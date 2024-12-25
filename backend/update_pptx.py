import json
import sys
from pptx import Presentation

# Debugging: Print received arguments
print("Arguments:", sys.argv)

if len(sys.argv) < 4:
    print("Error: Not enough arguments provided!")
    sys.exit(1)

template_path = sys.argv[1]
slides_json = sys.argv[2]
output_path = sys.argv[3]

# Parse JSON data
try:
    slides = json.loads(slides_json)
    print("Slides data parsed:", slides)
except json.JSONDecodeError as e:
    print("JSON parsing error:", e)
    sys.exit(1)

# Load presentation template
presentation = Presentation(template_path)

# Add slides to the presentation
for slide_data in slides:
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Use a blank layout
    title = slide.shapes.title
    title.text = slide_data["title"]

    textbox = slide.shapes.add_textbox(0, 100, 500, 300)
    textbox.text = slide_data["paragraph"]

# Save the updated presentation
presentation.save(output_path)
print("Presentation updated successfully:", output_path)
