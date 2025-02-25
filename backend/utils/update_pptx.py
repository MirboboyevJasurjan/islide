import sys
import json
import os
from pptx import Presentation

def replace_text_in_shape(shape, replacements):
    if not hasattr(shape, "text_frame"):
        return
    text_frame = shape.text_frame
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            for placeholder, new_text in replacements.items():
                if placeholder in run.text:
                    run.text = run.text.replace(placeholder, new_text)

def update_pptx(template_path, output_path, slides_data):
    prs = Presentation(template_path)
    
    # Iterate over each slide in the template
    for i, slide in enumerate(prs.slides):
        if i < len(slides_data):
            replacements = {
                'header': slides_data[i].get('header', ''),
                'title': slides_data[i].get('title', ''),
                'paragraph': slides_data[i].get('paragraph', '')
            }
            for shape in slide.shapes:
                replace_text_in_shape(shape, replacements)
    
    prs.save(output_path)
    print(f"PPTX file generated: {output_path}")

if __name__ == "__main__":
    try:
        template_path = sys.argv[1]
        slides_json = json.loads(sys.argv[2])
        output_path = sys.argv[3]  # Make sure it's inside /generated
        
        # Ensure /generated directory exists
        generated_dir = os.path.dirname(output_path)
        if not os.path.exists(generated_dir):
            os.makedirs(generated_dir)
        
        update_pptx(template_path, output_path, slides_json)
    except Exception as e:
        print("Error:", e)
        sys.exit(1)